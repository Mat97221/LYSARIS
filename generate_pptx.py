from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Color palette
BG_IVORY = "F4F0E8"
BG_DARK = "0C0F0A"
BG_IVORY2 = "EDE8DC"
BG_DARK2 = "1A1F17"
SAGE = "5C7A4E"
GOLD = "B8922A"
DARK_INK = "0C0F0A"
LIGHT_TEXT = "F4F0E8"
MUTED = "7A8070"
WHITE = "FFFFFF"

def hex_to_rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_rect(slide, left, top, width, height, fill_hex=None, line_hex=None, line_width=Pt(1)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    fill = shape.fill
    if fill_hex:
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(fill_hex)
    else:
        fill.background()
    line = shape.line
    if line_hex:
        line.color.rgb = hex_to_rgb(line_hex)
        line.width = line_width
    else:
        line.fill.background()
    return shape

def add_oval(slide, left, top, width, height, fill_hex, line_hex=None):
    shape = slide.shapes.add_shape(
        9,  # oval
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = hex_to_rgb(line_hex)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, left, top, width, height,
                font_name="Calibri", font_size=12, bold=False, italic=False,
                color_hex=DARK_INK, align=PP_ALIGN.LEFT, word_wrap=True,
                vertical_anchor=None):
    from pptx.enum.text import MSO_ANCHOR
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    if vertical_anchor:
        tf.vertical_anchor = vertical_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = hex_to_rgb(color_hex)
    return txBox

def add_textbox_multirun(slide, runs, left, top, width, height,
                          align=PP_ALIGN.LEFT, word_wrap=True):
    """runs: list of dicts with keys: text, font_size, bold, italic, color_hex, font_name"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    for r in runs:
        run = p.add_run()
        run.text = r.get("text","")
        run.font.name = r.get("font_name","Calibri")
        run.font.size = Pt(r.get("font_size",12))
        run.font.bold = r.get("bold", False)
        run.font.italic = r.get("italic", False)
        run.font.color.rgb = hex_to_rgb(r.get("color_hex", DARK_INK))
    return txBox

def set_slide_bg(slide, color_hex):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)

def add_para_to_tf(tf, text, font_size=12, bold=False, italic=False,
                   color_hex=DARK_INK, align=PP_ALIGN.LEFT, font_name="Calibri",
                   space_before=0):
    from pptx.util import Pt as PT
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = PT(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = hex_to_rgb(color_hex)
    return p

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# ─── SLIDE 1 ───────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)

# Kicker
add_textbox(slide, "Présentation à la CCI de l'Allier",
            0.6, 0.5, 12, 0.35, font_size=14, color_hex=GOLD)

# Title
add_textbox(slide, "L'infrastructure numérique des marchés de gros alimentaires français",
            0.6, 1.05, 12.0, 1.8, font_size=44, bold=True, color_hex=LIGHT_TEXT, word_wrap=True)

# Subtitle
add_textbox(slide,
    "LYSARIS construit l'outil de référence de l'information de marché pour les Marchés d'Intérêt National et les circuits de gros — un écosystème d'environ 40 milliards d'euros, structurellement sous-équipé en outils numériques.",
    0.6, 3.1, 12.0, 1.6, font_size=18, color_hex=LIGHT_TEXT, word_wrap=True)

# Pills at bottom
pills = [
    "Phase pilote — partenaires sélectionnés",
    "Fondée en 2026 — trois associés",
    "SAS en cours de finalisation",
]
pill_widths = [3.2, 2.8, 2.8]
pill_starts = [0.6, 4.05, 7.1]
for i, (pill_text, pw, ps) in enumerate(zip(pills, pill_widths, pill_starts)):
    add_rect(slide, ps, 6.5, pw, 0.45, fill_hex=None, line_hex=GOLD, line_width=Pt(1.2))
    add_textbox(slide, pill_text, ps+0.1, 6.52, pw-0.2, 0.41,
                font_size=11, color_hex=LIGHT_TEXT, align=PP_ALIGN.CENTER)

# ─── SLIDE 2 ───────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_IVORY)

add_textbox(slide, "ACCROCHE TERRITORIALE",
            0.6, 0.35, 12, 0.3, font_size=12, color_hex=SAGE)
add_textbox(slide, "Un projet né au cœur d'une région agricole et agroalimentaire",
            0.6, 0.75, 12, 0.8, font_size=32, bold=True, color_hex=DARK_INK, word_wrap=True)
add_textbox(slide,
    "L'Allier et la région Auvergne-Rhône-Alpes concentrent un tissu dense de producteurs, de grossistes, de coopératives et de plateformes de distribution alimentaire. C'est précisément pour ces acteurs — ceux qui font tourner les circuits de gros au quotidien — que LYSARIS développe ses outils.",
    0.6, 1.75, 12, 1.2, font_size=14, color_hex=DARK_INK, word_wrap=True)

# Quote block with gold left border
add_rect(slide, 0.6, 3.1, 0.07, 1.0, fill_hex=GOLD)
add_textbox(slide,
    "Notre conviction : la modernisation des circuits de gros alimentaires ne viendra pas des géants du numérique, mais d'acteurs de terrain qui comprennent comment fonctionne réellement un carreau de marché à cinq heures du matin.",
    0.8, 3.1, 11.6, 1.0, font_size=13, italic=True, color_hex=MUTED, word_wrap=True)

# KPI boxes
kpis = [
    ("~40 Md€", "Écosystème des marchés de gros français"),
    ("4", "Briques produit séquencées"),
    ("3", "Associés fondateurs complémentaires"),
]
kpi_x = [0.6, 4.77, 8.94]
for i, ((val, label), kx) in enumerate(zip(kpis, kpi_x)):
    add_rect(slide, kx, 4.35, 3.8, 2.6, fill_hex=BG_DARK)
    add_textbox(slide, val, kx+0.15, 4.55, 3.5, 0.9,
                font_size=34, bold=True, color_hex=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, kx+0.15, 5.55, 3.5, 1.2,
                font_size=12, color_hex=LIGHT_TEXT, align=PP_ALIGN.CENTER, word_wrap=True)

# ─── SLIDE 3 ───────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_IVORY2)

add_textbox(slide, "QUI NOUS SOMMES",
            0.6, 0.3, 12, 0.3, font_size=12, color_hex=SAGE)
add_textbox(slide, "Trois associés, trois métiers complémentaires",
            0.6, 0.68, 12, 0.7, font_size=32, bold=True, color_hex=DARK_INK)

cards = [
    ("Alex", "Président",
     "Direction commerciale et relations terrain. Profil technico-commercial issu de l'industrie, en charge du démarchage des marchés et des partenariats pilotes."),
    ("Stan", "Directeur Technique",
     "Architecture et développement de la plateforme. Profil ingénierie logicielle et cybersécurité — une exigence forte de fiabilité et de protection des données dès la conception."),
    ("Mat", "Directeur Général · Directeur Financier",
     "Structuration juridique, financière et conformité, notamment la protection des données personnelles (RGPD) des professionnels utilisateurs."),
]
card_x = [0.6, 4.77, 8.94]
for i, ((name, role, desc), cx) in enumerate(zip(cards, card_x)):
    add_rect(slide, cx, 1.6, 3.8, 5.4, fill_hex=WHITE)
    add_textbox(slide, name, cx+0.2, 1.85, 3.4, 0.7,
                font_size=26, bold=True, color_hex=DARK_INK)
    add_textbox(slide, role, cx+0.2, 2.65, 3.4, 0.55,
                font_size=13, color_hex=SAGE)
    add_textbox(slide, desc, cx+0.2, 3.3, 3.4, 3.4,
                font_size=13, color_hex=DARK_INK, word_wrap=True)

# ─── SLIDE 4 ───────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK2)

add_textbox(slide, "LE CONSTAT",
            0.6, 0.3, 12, 0.3, font_size=12, color_hex=GOLD)
add_textbox(slide, "Un écosystème majeur qui travaille sans visibilité partagée",
            0.6, 0.7, 12, 0.8, font_size=32, bold=True, color_hex=LIGHT_TEXT, word_wrap=True)
add_textbox(slide,
    "Sur les marchés de gros, grossistes, producteurs et acheteurs professionnels travaillent chaque jour sans information fiable et partagée sur les cours, les disponibilités et les tendances. L'information circule par téléphone, papier et habitude.",
    0.6, 1.75, 12, 1.1, font_size=14, color_hex=LIGHT_TEXT, word_wrap=True)

# Quote
add_rect(slide, 0.6, 3.0, 0.07, 1.1, fill_hex=GOLD)
add_textbox(slide,
    "Les conséquences sont concrètes : prix décorrélés du marché, ruptures non anticipées, opportunités manquées — pour des entreprises souvent familiales, dont la marge se joue à quelques pourcents.",
    0.8, 3.0, 11.6, 1.1, font_size=13, italic=True, color_hex=GOLD, word_wrap=True)

add_textbox(slide,
    "Les outils numériques généralistes ne répondent pas aux codes de ce milieu : horaires décalés, relations de confiance construites sur le carreau, usages mobiles simples.",
    0.6, 4.3, 12, 0.9, font_size=12, color_hex=MUTED, word_wrap=True)

# ─── SLIDE 5 ───────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_IVORY)

add_textbox(slide, "NOTRE RÉPONSE",
            0.6, 0.25, 12, 0.3, font_size=12, color_hex=SAGE)
add_textbox(slide, "Quatre briques, déployées dans l'ordre de la confiance",
            0.6, 0.62, 12, 0.7, font_size=32, bold=True, color_hex=DARK_INK)
add_textbox(slide,
    "Notre positionnement : devenir d'abord la référence de l'information de marché — « le Bloomberg des marchés de gros » — avant d'outiller la transaction.",
    0.6, 1.42, 12, 0.55, font_size=13, color_hex=MUTED, word_wrap=True)

# 4 product rows
rows = [
    ("IRIS",   SAGE,    WHITE,    "Observatoire de marché en temps réel",
     "Cours du jour, comparaisons anonymisées, historiques de prix, alertes.",
     "En cours de déploiement", GOLD),
    ("ORACLE", GOLD,    DARK_INK, "Couche d'anticipation",
     "Alertes sur prix anormaux, anticipation des pics saisonniers à 7-14 jours, substitutions.",
     "2e semestre", GOLD),
    ("NEXUS",  MUTED,   WHITE,    "Déploiement multi-sites",
     "La même infrastructure répliquée marché par marché, données cloisonnées par site.",
     "Horizon 12 mois", MUTED),
    ("AGORA",  BG_DARK, LIGHT_TEXT,"Facilitation des transactions",
     "Mise en relation et passation de commandes entre professionnels.",
     "Au-delà de 12 mois", MUTED),
]
row_bg_colors = [BG_IVORY, BG_IVORY2, BG_IVORY, BG_IVORY2]
row_y_start = 2.1
row_h = 1.15
for i, (badge, badge_bg, badge_fg, title, desc, status, status_color) in enumerate(rows):
    ry = row_y_start + i * row_h
    add_rect(slide, 0.4, ry, 12.5, row_h - 0.05, fill_hex=row_bg_colors[i])
    # Badge
    add_rect(slide, 0.55, ry + 0.2, 1.1, 0.55, fill_hex=badge_bg)
    add_textbox(slide, badge, 0.55, ry + 0.22, 1.1, 0.51,
                font_size=11, bold=True, color_hex=badge_fg, align=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, title, 1.8, ry + 0.1, 4.5, 0.45,
                font_size=13, bold=True, color_hex=DARK_INK)
    # Desc
    add_textbox(slide, desc, 1.8, ry + 0.57, 6.5, 0.55,
                font_size=11, color_hex=DARK_INK, word_wrap=True)
    # Status
    add_textbox(slide, status, 9.6, ry + 0.3, 3.1, 0.45,
                font_size=11, color_hex=status_color, align=PP_ALIGN.RIGHT)

# ─── SLIDE 6 ───────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_IVORY2)

add_textbox(slide, "OÙ NOUS ALLONS",
            0.6, 0.3, 12, 0.3, font_size=12, color_hex=SAGE)
add_textbox(slide, "Une feuille de route prudente, jalonnée par les résultats",
            0.6, 0.68, 12, 0.8, font_size=32, bold=True, color_hex=DARK_INK, word_wrap=True)

timeline_items = [
    (GOLD,  True,  "Réalisé 2026",            DARK_INK, "Fondation et préparation pilote",
     "Constitution de l'équipe, conception IRIS, cadrage juridique."),
    (SAGE,  True,  "Court terme — en cours",   SAGE,     "Premiers partenariats pilotes",
     "Signature partenariats gratuits fournisseurs, base de données cours propriétaire, conformité RGPD."),
    (MUTED, False, "6–12 mois",                DARK_INK, "Premiers revenus récurrents",
     "Abonnement acheteurs, levée de fonds amorçage, JEI, CIR."),
    (MUTED, False, "12–36 mois",               DARK_INK, "Déploiement national",
     "Réplication sur plusieurs marchés, partenariats FranceAgriMer, indices de prix."),
]
tl_y = 1.75
for i, (dot_color, is_filled, period, period_color, title, desc) in enumerate(timeline_items):
    item_y = tl_y + i * 1.32
    # Dot
    if is_filled:
        add_oval(slide, 0.55, item_y + 0.12, 0.22, 0.22, fill_hex=dot_color)
    else:
        dot = add_oval(slide, 0.55, item_y + 0.12, 0.22, 0.22, fill_hex=BG_IVORY2, line_hex=dot_color)
        dot.line.color.rgb = hex_to_rgb(dot_color)
        dot.line.width = Pt(1.5)
    # Vertical connector (except last)
    if i < 3:
        add_rect(slide, 0.645, item_y + 0.35, 0.03, 1.0, fill_hex=MUTED)
    # Period
    add_textbox(slide, period, 0.95, item_y, 3.0, 0.35,
                font_size=13, bold=is_filled, color_hex=period_color)
    # Title
    add_textbox(slide, title, 0.95, item_y + 0.37, 11.5, 0.38,
                font_size=14, bold=True, color_hex=DARK_INK)
    # Desc
    add_textbox(slide, desc, 0.95, item_y + 0.77, 11.5, 0.45,
                font_size=12, color_hex=MUTED, word_wrap=True)

# ─── SLIDE 7 ───────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK2)

add_textbox(slide, "RETOMBÉES POUR LE TERRITOIRE",
            0.6, 0.3, 12, 0.3, font_size=12, color_hex=GOLD)
add_textbox(slide, "Ce que le projet apporte localement",
            0.6, 0.7, 12, 0.7, font_size=32, bold=True, color_hex=LIGHT_TEXT)

s7_cards = [
    ("Visibilité économique",
     "Les producteurs et grossistes locaux accèdent gratuitement à une lecture fiable de leur marché."),
    ("Modernisation des circuits de gros",
     "Une infrastructure numérique pensée pour les usages réels du métier."),
    ("Emplois qualifiés en région",
     "Des postes qualifiés — données, développement commercial — ancrés en région."),
]
for i, (title, desc) in enumerate(s7_cards):
    cx = 0.6 + i * 4.17
    add_rect(slide, cx, 1.65, 3.9, 5.3, fill_hex="222820")
    add_rect(slide, cx, 1.65, 3.9, 0.09, fill_hex=SAGE)
    add_textbox(slide, title, cx+0.2, 1.9, 3.5, 0.75,
                font_size=16, bold=True, color_hex=LIGHT_TEXT, word_wrap=True)
    add_textbox(slide, desc, cx+0.2, 2.8, 3.5, 3.8,
                font_size=13, color_hex=LIGHT_TEXT, word_wrap=True)

# ─── SLIDE 8 ───────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_IVORY)

add_textbox(slide, "NOTRE DEMANDE À LA CCI DE L'ALLIER",
            0.6, 0.3, 12, 0.3, font_size=12, color_hex=SAGE)
add_textbox(slide, "Trois accompagnements concrets",
            0.6, 0.68, 12, 0.7, font_size=32, bold=True, color_hex=DARK_INK)

asks = [
    ("1", "Mise en relation avec les acteurs agroalimentaires du département",
     "Grossistes, coopératives, plateformes de distribution et logistique : une introduction qualifiée par la CCI vaut des mois de prospection à froid."),
    ("2", "Orientation vers les dispositifs régionaux d'appui à l'innovation",
     "Région Auvergne-Rhône-Alpes, Bpifrance, réseau French Tech local."),
    ("3", "Conseil sur l'ancrage territorial du projet",
     "Implantation, recrutements futurs, participation aux filières agroalimentaires locales."),
]
ask_y = 1.65
for i, (num, title, desc) in enumerate(asks):
    ry = ask_y + i * 1.65
    add_rect(slide, 0.45, ry, 12.4, 1.5, fill_hex=BG_IVORY2)
    add_textbox(slide, num, 0.6, ry + 0.25, 0.6, 0.9,
                font_size=36, bold=True, color_hex=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, title, 1.35, ry + 0.15, 11.1, 0.5,
                font_size=14, bold=True, color_hex=DARK_INK, word_wrap=True)
    add_textbox(slide, desc, 1.35, ry + 0.68, 11.1, 0.75,
                font_size=13, color_hex=DARK_INK, word_wrap=True)

# ─── SLIDE 9 ───────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)

add_textbox(slide, "LYSARIS",
            0.5, 1.15, 12.33, 1.5,
            font_size=60, bold=True, color_hex=LIGHT_TEXT, align=PP_ALIGN.CENTER)
add_textbox(slide,
    "L'infrastructure numérique des marchés de gros alimentaires français",
    0.5, 2.7, 12.33, 0.6,
    font_size=16, color_hex=GOLD, align=PP_ALIGN.CENTER)
add_textbox(slide,
    "Jeune entreprise fondée en 2026 — phase pilote avec partenaires sélectionnés.",
    0.5, 3.45, 12.33, 0.5,
    font_size=13, color_hex=LIGHT_TEXT, align=PP_ALIGN.CENTER)
add_textbox(slide, "MAT — Directeur Général",
            0.5, 4.3, 12.33, 0.45,
            font_size=14, color_hex=LIGHT_TEXT, align=PP_ALIGN.CENTER)
add_textbox(slide, "contact@lysaris.fr",
            0.5, 4.85, 12.33, 0.45,
            font_size=14, color_hex=GOLD, align=PP_ALIGN.CENTER)
add_textbox(slide, "SAS LYSARIS — STRUCTURE EN COURS DE FINALISATION",
            0.5, 6.9, 12.33, 0.4,
            font_size=10, color_hex=MUTED, align=PP_ALIGN.CENTER)

# Save
output_path = "/home/user/LYSARIS/LYSARIS_CCI_Allier_2026.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
